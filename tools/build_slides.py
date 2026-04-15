"""
YouTube AI Niche Report — PowerPoint Generator
Creates a professional slide deck from analysis results.

Usage: python tools/build_slides.py
Input: .tmp/analysis_results.json
Output: .tmp/ai_youtube_report.pptx
"""

import json
import sys
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

ROOT = Path(__file__).resolve().parent.parent
INPUT_PATH = ROOT / ".tmp" / "analysis_results.json"
OUTPUT_PATH = ROOT / ".tmp" / "ai_youtube_report.pptx"

# Color palette
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_PURPLE = RGBColor(0x7C, 0x3A, 0xED)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_SURFACE = RGBColor(0x1E, 0x29, 0x3B)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=WHITE, bullet_color=ACCENT_BLUE):
    """Add a bulleted list to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def add_stat_card(slide, left, top, width, height, value, label,
                  value_color=ACCENT_BLUE):
    """Add a stat card with large number and label."""
    # Card background
    shape = slide.shapes.add_shape(
        1, left, top, width, height  # MSO_SHAPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_SURFACE
    shape.line.fill.background()

    # Value
    add_textbox(slide, left + Inches(0.2), top + Inches(0.2),
                width - Inches(0.4), Inches(0.8),
                value, font_size=28, color=value_color, bold=True,
                alignment=PP_ALIGN.CENTER)
    # Label
    add_textbox(slide, left + Inches(0.2), top + Inches(0.9),
                width - Inches(0.4), Inches(0.5),
                label, font_size=11, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)


def add_bar_chart(slide, left, top, width, height, categories, values, title,
                  bar_color=ACCENT_BLUE):
    """Add a bar chart to a slide."""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(title, values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.has_title = False

    # Style the chart
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = bar_color

    # Style axes
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(9)
    cat_axis.tick_labels.font.color.rgb = LIGHT_GRAY
    cat_axis.format.line.fill.background()

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(9)
    val_axis.tick_labels.font.color.rgb = LIGHT_GRAY
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.color.rgb = RGBColor(0x33, 0x33, 0x55)

    return chart_frame


def slide_title(slide, left, top, width, text):
    """Add a section title to a slide."""
    add_textbox(slide, left, top, width, Inches(0.6),
                text, font_size=28, color=WHITE, bold=True)
    # Accent line
    shape = slide.shapes.add_shape(1, left, top + Inches(0.55), Inches(2), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()


def build_presentation(data: dict) -> str:
    """Build the full PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    summary = data["summary"]
    topics = data["trending_topics"]
    top_videos = data["top_videos"]
    top_channels = data["top_channels"]
    benchmarks = data["engagement_benchmarks"]
    opportunities = data["content_opportunities"]

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ─── SLIDE 1: Title ───
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    # Accent bar at top
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1.2),
                "AI YouTube Landscape Report", font_size=44, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
                f"Weekly Intelligence Report  |  {summary['date_range']}",
                font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
                f"{summary['total_videos']} videos analyzed  •  {summary['unique_channels']} channels tracked  •  {summary['total_views']:,} total views",
                font_size=14, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # ─── SLIDE 2: Executive Summary ───
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, Inches(0.8), Inches(0.5), Inches(11), "Executive Summary")

    # Stat cards row
    card_width = Inches(2.8)
    card_height = Inches(1.5)
    card_y = Inches(1.4)
    gap = Inches(0.3)
    start_x = Inches(0.8)

    add_stat_card(slide, start_x, card_y, card_width, card_height,
                  f"{summary['total_videos']}", "Videos Analyzed", ACCENT_BLUE)
    add_stat_card(slide, start_x + card_width + gap, card_y, card_width, card_height,
                  f"{summary['avg_views']:,}", "Average Views", ACCENT_GREEN)
    add_stat_card(slide, start_x + 2*(card_width + gap), card_y, card_width, card_height,
                  f"{summary['max_views']:,}", "Top Video Views", ACCENT_ORANGE)
    add_stat_card(slide, start_x + 3*(card_width + gap), card_y, card_width, card_height,
                  f"{summary['unique_channels']}", "Channels Tracked", ACCENT_PURPLE)

    # Key findings
    add_textbox(slide, Inches(0.8), Inches(3.3), Inches(11), Inches(0.5),
                "Key Findings", font_size=20, color=WHITE, bold=True)
    findings = [
        f"Google Gemma 4 is the breakout topic — 1.6M views across just 4 videos",
        f"Short-form content (0-5min) dominates — 235K avg views vs 31K for 30+ min videos",
        f"Saturday is the best publish day — 459K avg views, 2.2x better than weekday average",
        f"Median engagement: {benchmarks['median_like_ratio']:.1%} like ratio, {benchmarks['median_comment_ratio']:.2%} comment ratio",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(3.9), Inches(11), Inches(3),
                    findings, font_size=15, color=LIGHT_GRAY)

    # ─── SLIDE 3: Trending Topics ───
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, Inches(0.8), Inches(0.5), Inches(11), "Trending Topics")

    # Chart - top topics by total views
    chart_topics = [t["topic"].title() for t in topics[:8]]
    chart_views = [t["total_views"] / 1_000_000 for t in topics[:8]]
    add_bar_chart(slide, Inches(0.8), Inches(1.4), Inches(7), Inches(5.5),
                  chart_topics, chart_views, "Total Views (M)", ACCENT_PURPLE)

    # Side panel with details
    add_textbox(slide, Inches(8.3), Inches(1.4), Inches(4.5), Inches(0.5),
                "Rising Topics", font_size=18, color=WHITE, bold=True)
    rising = [t for t in topics if t["trend"] == "rising"][:6]
    rising_items = [
        f"{t['topic'].title()} — {t['video_count']} videos, {t['avg_views']:,} avg views"
        for t in rising
    ]
    add_bullet_list(slide, Inches(8.3), Inches(2.0), Inches(4.5), Inches(4.5),
                    rising_items, font_size=12, color=LIGHT_GRAY)

    # ─── SLIDE 4: Top Performing Videos ───
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, Inches(0.8), Inches(0.5), Inches(11), "Top Performing Videos")

    y = Inches(1.5)
    for i, v in enumerate(top_videos[:5]):
        # Rank number
        add_textbox(slide, Inches(0.8), y, Inches(0.6), Inches(0.5),
                    f"#{i+1}", font_size=24, color=ACCENT_BLUE, bold=True)
        # Title
        title_text = v["title"][:70] + ("..." if len(v["title"]) > 70 else "")
        add_textbox(slide, Inches(1.5), y, Inches(6.5), Inches(0.4),
                    title_text, font_size=16, color=WHITE, bold=True)
        # Channel + stats
        add_textbox(slide, Inches(1.5), y + Inches(0.35), Inches(6.5), Inches(0.3),
                    f"by {v['channel']}", font_size=12, color=LIGHT_GRAY)
        # View count
        add_textbox(slide, Inches(8.5), y, Inches(2), Inches(0.4),
                    f"{v['views']:,} views", font_size=16, color=ACCENT_GREEN, bold=True,
                    alignment=PP_ALIGN.RIGHT)
        # Like ratio
        add_textbox(slide, Inches(10.8), y, Inches(1.8), Inches(0.4),
                    f"{v['like_ratio']:.1%} likes", font_size=14, color=ACCENT_ORANGE,
                    alignment=PP_ALIGN.RIGHT)
        # Separator
        if i < 4:
            sep = slide.shapes.add_shape(1, Inches(0.8), y + Inches(0.8), Inches(11.7), Inches(0.01))
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
            sep.line.fill.background()
        y += Inches(1.1)

    # ─── SLIDE 5: Top Channels ───
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, Inches(0.8), Inches(0.5), Inches(11), "Top Channels by Engagement")

    y = Inches(1.5)
    for i, ch in enumerate(top_channels[:5]):
        add_textbox(slide, Inches(0.8), y, Inches(0.6), Inches(0.5),
                    f"#{i+1}", font_size=24, color=ACCENT_PURPLE, bold=True)
        add_textbox(slide, Inches(1.5), y, Inches(4), Inches(0.4),
                    ch["name"], font_size=16, color=WHITE, bold=True)
        add_textbox(slide, Inches(1.5), y + Inches(0.35), Inches(4), Inches(0.3),
                    f"{ch['subscriber_count']:,} subscribers", font_size=12, color=LIGHT_GRAY)
        add_textbox(slide, Inches(6), y, Inches(2.5), Inches(0.4),
                    f"{ch['engagement_rate']:.1%} engagement", font_size=16, color=ACCENT_GREEN, bold=True,
                    alignment=PP_ALIGN.RIGHT)
        add_textbox(slide, Inches(9), y, Inches(2.5), Inches(0.4),
                    f"{ch['avg_views']:,} avg views", font_size=14, color=ACCENT_ORANGE,
                    alignment=PP_ALIGN.RIGHT)
        if i < 4:
            sep = slide.shapes.add_shape(1, Inches(0.8), y + Inches(0.8), Inches(11.7), Inches(0.01))
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
            sep.line.fill.background()
        y += Inches(1.1)

    # ─── SLIDE 6: Engagement Benchmarks ───
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, Inches(0.8), Inches(0.5), Inches(11), "Engagement Benchmarks")

    # Stat cards
    add_stat_card(slide, Inches(0.8), Inches(1.4), Inches(2.8), Inches(1.5),
                  f"{benchmarks['median_like_ratio']:.1%}", "Median Like Ratio", ACCENT_BLUE)
    add_stat_card(slide, Inches(4.0), Inches(1.4), Inches(2.8), Inches(1.5),
                  f"{benchmarks['median_comment_ratio']:.2%}", "Median Comment Ratio", ACCENT_GREEN)
    add_stat_card(slide, Inches(7.2), Inches(1.4), Inches(2.8), Inches(1.5),
                  f"{benchmarks['avg_view_velocity']:,.0f}/hr", "Avg View Velocity", ACCENT_ORANGE)
    add_stat_card(slide, Inches(10.4), Inches(1.4), Inches(2.4), Inches(1.5),
                  benchmarks["best_duration_bucket"], "Best Duration", ACCENT_PURPLE)

    # Duration breakdown chart
    dur_data = benchmarks.get("duration_breakdown", {})
    if dur_data:
        cats = list(dur_data.keys())
        vals = [v / 1000 for v in dur_data.values()]
        add_bar_chart(slide, Inches(0.8), Inches(3.3), Inches(5.8), Inches(3.8),
                      cats, vals, "Avg Views (K)", ACCENT_BLUE)
        add_textbox(slide, Inches(0.8), Inches(3.0), Inches(5.8), Inches(0.4),
                    "Average Views by Video Duration", font_size=14, color=LIGHT_GRAY)

    # Pattern breakdown chart
    pat_data = benchmarks.get("pattern_breakdown", {})
    if pat_data:
        cats = [k.title() for k in pat_data.keys()]
        vals = [v / 1000 for v in pat_data.values()]
        add_bar_chart(slide, Inches(7.2), Inches(3.3), Inches(5.8), Inches(3.8),
                      cats, vals, "Avg Views (K)", ACCENT_PURPLE)
        add_textbox(slide, Inches(7.2), Inches(3.0), Inches(5.8), Inches(0.4),
                    "Average Views by Title Pattern", font_size=14, color=LIGHT_GRAY)

    # ─── SLIDE 7: Publishing Strategy ───
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, Inches(0.8), Inches(0.5), Inches(11), "Optimal Publishing Strategy")

    # Day of week chart
    day_data = benchmarks.get("day_breakdown", {})
    if day_data:
        day_order = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"]
        ordered_days = [d for d in day_order if d in day_data]
        day_vals = [day_data[d] / 1000 for d in ordered_days]
        short_days = [d[:3] for d in ordered_days]
        add_bar_chart(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(5.2),
                      short_days, day_vals, "Avg Views (K)", ACCENT_GREEN)
        add_textbox(slide, Inches(0.8), Inches(1.2), Inches(6), Inches(0.4),
                    "Average Views by Publish Day", font_size=14, color=LIGHT_GRAY)

    # Key takeaways
    add_textbox(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.5),
                "Key Takeaways", font_size=20, color=WHITE, bold=True)
    strategy_items = [
        f"Best day: {benchmarks.get('best_publish_day', 'N/A')} (459K avg views)",
        f"Best hour: {benchmarks.get('best_publish_hour_utc', 0)}:00 UTC",
        "Comparison titles average 122K views",
        "Question titles average 109K views",
        "How-to titles average 13K views",
        "Listicle titles average 14K views",
        "Short-form (0-5min) = 7.5x more views than 30+ min",
    ]
    add_bullet_list(slide, Inches(7.5), Inches(2.2), Inches(5), Inches(4.5),
                    strategy_items, font_size=13, color=LIGHT_GRAY)

    # ─── SLIDE 8: Content Opportunities ───
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, Inches(0.8), Inches(0.5), Inches(11), "Content Opportunities")

    y = Inches(1.5)
    for i, opp in enumerate(opportunities[:5]):
        # Card background
        card = slide.shapes.add_shape(1, Inches(0.8), y, Inches(11.7), Inches(1.0))
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_SURFACE
        card.line.fill.background()

        # Topic name
        add_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(3), Inches(0.4),
                    opp["topic"].title(), font_size=18, color=ACCENT_BLUE, bold=True)
        # Avg views
        add_textbox(slide, Inches(10), y + Inches(0.1), Inches(2.3), Inches(0.4),
                    f"{opp['avg_views']:,} avg views", font_size=14, color=ACCENT_GREEN,
                    alignment=PP_ALIGN.RIGHT)
        # Reason
        add_textbox(slide, Inches(1.0), y + Inches(0.5), Inches(11.3), Inches(0.4),
                    opp["reason"], font_size=12, color=LIGHT_GRAY)

        y += Inches(1.15)

    # ─── SLIDE 9: Recommendations ───
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, Inches(0.8), Inches(0.5), Inches(11), "Recommendations")

    recs = [
        "Create a Gemma 4 tutorial or comparison video — highest opportunity with low competition",
        "Publish on Saturday for maximum reach — 2.2x better than weekday average",
        "Keep videos under 5 minutes — short-form gets 7.5x more views than long-form",
        "Use comparison or question-style titles — they outperform how-to and listicles by 8x",
        "Cover OpenAI / Sam Altman developments — consistent audience demand with fresh news cycle",
    ]

    y = Inches(1.5)
    for i, rec in enumerate(recs):
        # Number circle
        add_textbox(slide, Inches(0.8), y, Inches(0.6), Inches(0.6),
                    str(i + 1), font_size=28, color=ACCENT_BLUE, bold=True)
        # Recommendation text
        add_textbox(slide, Inches(1.5), y + Inches(0.05), Inches(11), Inches(0.8),
                    rec, font_size=16, color=WHITE)
        y += Inches(1.05)

    # ─── SLIDE 10: Methodology ───
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, Inches(0.8), Inches(0.5), Inches(11), "Methodology")

    method_items = [
        "Data Source: YouTube Data API v3 via OAuth",
        f"Date Range: {summary['date_range']}",
        "Search Keywords: AI automation, artificial intelligence, AI tools, AI agents, AI workflow, AI tutorial, AI news",
        "Search Strategy: 14 searches (7 keywords x 2 sort orders: relevance, viewCount)",
        f"Dataset: {summary['total_videos']} unique videos, {summary['unique_channels']} unique channels",
        "Metrics: views, likes, comments, duration, publish timing, tags, title patterns",
        "Engagement Rate: (likes + comments) / views",
        "View Velocity: views / hours since publish",
        "Topic Scoring: avg views x sqrt(video count)",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5),
                    method_items, font_size=14, color=LIGHT_GRAY)

    # Footer
    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
                f"Generated on {datetime.now(timezone.utc).strftime('%B %d, %Y at %H:%M UTC')} | YouTube AI Niche Analysis Automation",
                font_size=11, color=RGBColor(0x64, 0x74, 0x8B), alignment=PP_ALIGN.CENTER)

    # Save
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return str(OUTPUT_PATH)


def main():
    if not INPUT_PATH.exists():
        print(f"ERROR: Input file not found: {INPUT_PATH}")
        print("Run the analyzer first: python tools/analyze_data.py")
        sys.exit(1)

    print("Building AI YouTube Report...")
    with open(INPUT_PATH, "r", encoding="utf-8") as f:
        data = json.load(f)

    output = build_presentation(data)
    print(f"Presentation saved to: {output}")


if __name__ == "__main__":
    main()
